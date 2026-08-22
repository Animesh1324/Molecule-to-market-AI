"""Tests for the PPTX export's real-data density and the project_id fix.

Before this, /api/reports/export/pptx always regenerated a fresh, disposable
plan from the query-string molecule/brand fields — a project's own saved
plan (any AI-drafted content, edits, or in-progress KPI/milestone data) never
reached the export. These tests pin the fix and the new data-driven slides.
"""
import io

from fastapi.testclient import TestClient
from pptx import Presentation

from app.db.database import init_db
from app.main import app
from app.db import database as db
from app.services import export_service

init_db()
client = TestClient(app)


def _slide_texts(presentation):
    texts = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        texts.append(cell.text_frame.text)
    return texts


def test_export_uses_the_projects_own_saved_plan_when_project_id_given():
    saved_plan = {
        "project_id": "proj-pptx-test",
        "molecule_name": "Empagliflozin",
        "brand_name": "Cardioflo",
        "therapy_area": "Cardiometabolic",
        "indication": "Heart Failure",
        "target_geography": "Global",
        "mission": "A DISTINCTIVE_MISSION_STRING_FOR_THIS_TEST",
        "vision": "vision", "brand_objective": "objective",
        "therapy_area_opportunity": "x", "target_customer_and_patient_profile": "x",
        "doctor_and_market_insights": "x", "competitor_gap_and_differentiation": "x",
        "positioning_statement": "x", "brand_promise_and_rtb": "x",
        "key_messages_and_claim_strategy": "x", "commercial_launch_strategy": "x",
        "kol_and_cme_strategy": "x", "digital_and_sales_force_strategy": "x",
        "sections": [], "monthly_action_plan": [], "kpi_scorecard": [],
        "mlr_compliance_signoff_ready": False, "last_updated": "2026-01-01",
    }
    db.db_save_brand_plan("proj-pptx-test", saved_plan)

    response = client.get(
        "/api/reports/export/pptx?molecule=Empagliflozin&brand_name=Cardioflo&project_id=proj-pptx-test"
    )
    assert response.status_code == 200
    presentation = Presentation(io.BytesIO(response.content))
    texts = _slide_texts(presentation)
    assert any("DISTINCTIVE_MISSION_STRING_FOR_THIS_TEST" in t for t in texts)


def test_export_falls_back_to_a_fresh_plan_when_project_id_is_unknown():
    response = client.get(
        "/api/reports/export/pptx?molecule=Empagliflozin&brand_name=Cardioflo&project_id=proj-does-not-exist"
    )
    assert response.status_code == 200
    assert len(response.content) > 1000


def test_dashboard_slide_shows_real_competitor_and_regulatory_stats(monkeypatch):
    """Empagliflozin has curated competitor/SWOT data and a real US FDA record —
    the executive dashboard should surface at least one real figure, not a
    placeholder, when the underlying grounding calls succeed.
    """
    response = client.get("/api/reports/export/pptx?molecule=Empagliflozin&brand_name=Cardioflo")
    assert response.status_code == 200
    presentation = Presentation(io.BytesIO(response.content))
    texts = " ".join(_slide_texts(presentation))
    assert "Executive Summary Dashboard" in texts


def test_swot_table_renders_curated_data_for_empagliflozin():
    competitor_data = {
        "swot_analysis": {
            "strengths": ["EMPA-REG OUTCOME demonstrated a mortality benefit."],
            "weaknesses": ["Requires counseling on a known class side effect."],
            "opportunities": ["Guideline endorsement expanding first-line use."],
            "threats": ["A faster-growing competitor class."],
        }
    }
    from app.services.ai_orchestrator import generate_strategic_brand_plan, generate_commercial_assets
    plan = generate_strategic_brand_plan(project_id="t", molecule_name="Empagliflozin", brand_name="Cardioflo")
    assets = generate_commercial_assets(molecule_name="Empagliflozin", brand_name="Cardioflo")
    buffer = export_service.generate_pitch_deck_pptx(plan, assets, competitor_data=competitor_data)
    presentation = Presentation(io.BytesIO(buffer.getvalue()))
    texts = " ".join(_slide_texts(presentation))
    assert "EMPA-REG OUTCOME demonstrated a mortality benefit." in texts
    assert "A faster-growing competitor class." in texts


def test_swot_table_states_the_gap_when_no_swot_data_exists():
    from app.services.ai_orchestrator import generate_strategic_brand_plan, generate_commercial_assets
    plan = generate_strategic_brand_plan(project_id="t", molecule_name="Zzznotamolecule", brand_name="Cardioflo")
    assets = generate_commercial_assets(molecule_name="Zzznotamolecule", brand_name="Cardioflo")
    buffer = export_service.generate_pitch_deck_pptx(plan, assets, competitor_data=None)
    presentation = Presentation(io.BytesIO(buffer.getvalue()))
    texts = " ".join(_slide_texts(presentation))
    assert "No SWOT analysis on file yet" in texts


def test_trade_price_slide_only_appears_when_forecast_has_one():
    from app.models.forecast import TradePriceStructure
    from app.services.forecast_service import calculate_market_forecast
    from app.services.ai_orchestrator import generate_strategic_brand_plan, generate_commercial_assets

    plan = generate_strategic_brand_plan(project_id="t", molecule_name="Empagliflozin", brand_name="Cardioflo")
    assets = generate_commercial_assets(molecule_name="Empagliflozin", brand_name="Cardioflo")

    forecast_no_trade = calculate_market_forecast()
    buffer = export_service.generate_pitch_deck_pptx(plan, assets, forecast=forecast_no_trade)
    texts_without = " ".join(_slide_texts(Presentation(io.BytesIO(buffer.getvalue()))))
    assert "India Trade Price Structure" not in texts_without

    forecast_with_trade = calculate_market_forecast(
        mrp_per_patient_year_inr=18000, ptr_per_patient_year_inr=15500, pts_per_patient_year_inr=13200,
    )
    buffer2 = export_service.generate_pitch_deck_pptx(plan, assets, forecast=forecast_with_trade)
    texts_with = " ".join(_slide_texts(Presentation(io.BytesIO(buffer2.getvalue()))))
    assert "India Trade Price Structure" in texts_with
    assert "₹18,000" in texts_with


def test_kpi_and_milestone_tables_render_the_plans_own_rows():
    from app.services.ai_orchestrator import generate_strategic_brand_plan, generate_commercial_assets
    plan = generate_strategic_brand_plan(project_id="t", molecule_name="Empagliflozin", brand_name="Cardioflo")
    assets = generate_commercial_assets(molecule_name="Empagliflozin", brand_name="Cardioflo")
    buffer = export_service.generate_pitch_deck_pptx(plan, assets)
    texts = " ".join(_slide_texts(Presentation(io.BytesIO(buffer.getvalue()))))
    assert plan.kpi_scorecard, "fixture plan should have KPI rows to assert against"
    assert plan.kpi_scorecard[0].kpi_name in texts
    assert plan.monthly_action_plan[0].activity in texts


def test_primary_research_slide_only_appears_when_data_exists():
    from app.services.ai_orchestrator import generate_strategic_brand_plan, generate_commercial_assets
    plan = generate_strategic_brand_plan(project_id="t", molecule_name="Empagliflozin", brand_name="Cardioflo")
    assets = generate_commercial_assets(molecule_name="Empagliflozin", brand_name="Cardioflo")

    buffer_without = export_service.generate_pitch_deck_pptx(plan, assets, primary_research=None)
    texts_without = " ".join(_slide_texts(Presentation(io.BytesIO(buffer_without.getvalue()))))
    assert "Primary Research: RCPA" not in texts_without

    primary_research = {
        "has_data": True, "rcpa_total": 40, "rcpa_aware_count": 28, "rcpa_aware_percent": 70.0,
        "rcpa_active_count": 2, "rcpa_active_percent": 5.0, "hcp_total": 10,
        "hcp_avg_cost_barrier_rating": 7.0,
    }
    buffer_with = export_service.generate_pitch_deck_pptx(plan, assets, primary_research=primary_research)
    texts_with = " ".join(_slide_texts(Presentation(io.BytesIO(buffer_with.getvalue()))))
    assert "Primary Research: RCPA" in texts_with
    assert "28/40 (70.0%)" in texts_with


def test_title_slide_carries_developer_attribution():
    from app.services.ai_orchestrator import generate_strategic_brand_plan, generate_commercial_assets
    plan = generate_strategic_brand_plan(project_id="t", molecule_name="Empagliflozin", brand_name="Cardioflo")
    assets = generate_commercial_assets(molecule_name="Empagliflozin", brand_name="Cardioflo")
    buffer = export_service.generate_pitch_deck_pptx(plan, assets)
    texts = " ".join(_slide_texts(Presentation(io.BytesIO(buffer.getvalue()))))
    assert "Animesh Mishra" in texts

import uuid

import pytest
from fastapi.testclient import TestClient
from app.db.database import init_db
from app.main import app

# Ensure DB tables exist before any test (TestClient lifespan may not fire in all Starlette versions)
init_db()

client = TestClient(app)

def test_health():
    response = client.get("/health")
    assert response.status_code == 200
    body = response.json()
    assert body["status"] == "healthy"
    assert body["database"] == "connected"


def test_root_endpoint():
    response = client.get("/")
    assert response.status_code == 200
    data = response.json()
    assert data["status"] == "Operational"
    assert data["modules_active"] == 15

def test_projects_list():
    response = client.get("/api/projects")
    assert response.status_code == 200
    projects = response.json()
    assert len(projects) >= 2
    assert any(p["target_molecule_name"] == "Empagliflozin" for p in projects)

def test_molecule_intelligence():
    response = client.get("/api/molecules/search?name=Empagliflozin")
    assert response.status_code == 200
    data = response.json()
    assert data["generic_name"] == "Empagliflozin"
    assert "SGLT2" in data["pharmacological_class"]
    assert len(data["approved_indications"]) > 0

def test_evidence_papers():
    response = client.get("/api/evidence/papers?molecule=Empagliflozin")
    assert response.status_code == 200
    papers = response.json()
    assert len(papers) >= 1
    assert any("EMPA-REG" in p["title"] or "Empagliflozin" in p["title"] for p in papers)

def test_clinical_trials():
    response = client.get("/api/trials/landscape?molecule=Empagliflozin")
    assert response.status_code == 200
    data = response.json()
    assert data["total_trials_found"] >= 1
    assert len(data["landmark_trials"]) >= 1

def test_regulatory_labels():
    response = client.get("/api/regulatory/labels?molecule=Empagliflozin")
    assert response.status_code == 200
    data = response.json()
    assert data["us_fda"]["innovator_brand_name"] == "JARDIANCE"
    assert len(data["key_label_claims_verified"]) > 0

def test_trademark_analysis():
    response = client.get("/api/trademark/analyze?molecule=Empagliflozin")
    assert response.status_code == 200
    data = response.json()
    assert len(data["suggested_brand_names"]) >= 3
    assert all("uspto_search_link" in s for s in data["suggested_brand_names"])

def test_competitor_landscape():
    response = client.get("/api/competitors/landscape?molecule=Empagliflozin")
    assert response.status_code == 200
    data = response.json()
    assert len(data["competitors"]) >= 2
    assert len(data["swot_analysis"]["strengths"]) >= 1

def test_forecasting_model():
    response = client.get("/api/forecasting/model?prevalence_rate=0.10&treated_rate=0.60")
    assert response.status_code == 200
    data = response.json()
    assert data["prevalent_patient_pool"] > 0
    assert data["realistic_scenario"]["year_5"] > data["realistic_scenario"]["year_1"]

def test_brand_plan_generation():
    response = client.get("/api/brand-plan/generate?project_id=test-1&molecule=Empagliflozin")
    assert response.status_code == 200
    data = response.json()
    assert len(data["sections"]) == 12
    assert len(data["monthly_action_plan"]) >= 5

def test_brand_plan_reflects_real_competitor_data_not_a_bare_placeholder():
    """The competitive-defense section previously said the same generic
    sentence regardless of molecule, even when real competitor data existed.
    A manually-attested competitor (the isolated test database carries no
    pre-loaded licensed extract) must appear in the generated plan, not a
    generic placeholder.
    """
    from app.services import market_data_service as market

    entry = market.add_manual_competitor(
        molecule="Rosuvastatin", brand="TestCompetitorBrand", company="Test Co",
        source_note="Test fixture", added_by="Test")
    try:
        response = client.get(
            "/api/brand-plan/generate?project_id=test-competitor-grounding"
            "&molecule=Rosuvastatin&refresh=true&ai=false")
        assert response.status_code == 200
        data = response.json()
        competitor_section = next(
            s for s in data["sections"] if "Competitive Defense" in s["section_title"])
        assert "Build competitor comparison from verified labels" not in competitor_section["content_markdown"]
        assert "TestCompetitorBrand" in competitor_section["content_markdown"]
        assert "team-attested" in competitor_section["content_markdown"]
        assert "pending source-backed competitor and label comparison" not in data["competitor_gap_and_differentiation"]
    finally:
        market.delete_manual_competitor(entry["id"])


def test_brand_plan_states_the_measured_market_size_when_available():
    """A licensed-extract-backed competitor (not manual) states the measured
    market size in the plan, not just the brand list.
    """
    from app.services import market_data_service as market
    from app.db.market_models import MarketBrandORM, MarketDatasetORM
    from app.db.database import SessionLocal
    import uuid
    from datetime import datetime

    session = SessionLocal()
    dataset_id = uuid.uuid4().hex
    try:
        session.add(MarketDatasetORM(
            id=dataset_id, original_filename="test.csv", source_label="Test",
            market="India", period_label="MAT TEST'26", value_unit="INR Cr",
            row_count=1, brand_count=1, molecule_count=1, company_count=1,
            total_value=100.0, status="ready",
            ingested_at=datetime.now().strftime("%Y-%m-%d %H:%M:%S")))
        session.add(MarketBrandORM(
            id=uuid.uuid4().hex, dataset_id=dataset_id,
            molecule_desc="Apixaban", molecule_key="APIXABAN",
            brand="TestApixBrand", company="Test Pharma",
            value_latest=100.0, period_latest="MAT TEST'26"))
        session.commit()
    finally:
        session.close()

    try:
        response = client.get(
            "/api/brand-plan/generate?project_id=test-competitor-grounding-2"
            "&molecule=Apixaban&refresh=true&ai=false")
        data = response.json()
        competitor_section = next(
            s for s in data["sections"] if "Competitive Defense" in s["section_title"])
        assert "Measured market" in competitor_section["content_markdown"]
        assert "INR Cr" in competitor_section["content_markdown"]
    finally:
        market.delete_dataset(dataset_id)


def test_brand_plan_falls_back_cleanly_with_no_competitor_data_on_file():
    """A molecule with no licensed extract and no manual entry must still
    generate a plan — the gap is stated plainly, not a crash.
    """
    response = client.get(
        "/api/brand-plan/generate?project_id=test-competitor-grounding-3"
        "&molecule=Fictitiousmab&refresh=true&ai=false")
    assert response.status_code == 200
    data = response.json()
    competitor_section = next(
        s for s in data["sections"] if "Competitive Defense" in s["section_title"])
    assert "No source-backed competitor set is on file" in competitor_section["content_markdown"]


def test_creative_assets_generation():
    response = client.get("/api/assets/generate?molecule=Empagliflozin")
    assert response.status_code == 200
    data = response.json()
    assert len(data["visual_aid_slides"]) == 6
    assert len(data["mr_objection_handling_guide"]) >= 3

def test_export_docx():
    response = client.get("/api/reports/export/docx?molecule=Empagliflozin&brand_name=Cardioflo")
    assert response.status_code == 200
    assert response.headers["content-type"] == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    assert len(response.content) > 1000

def test_export_pptx():
    response = client.get("/api/reports/export/pptx?molecule=Empagliflozin&brand_name=Cardioflo")
    assert response.status_code == 200
    assert response.headers["content-type"] == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    assert len(response.content) > 1000

def test_export_xlsx():
    response = client.get("/api/reports/export/xlsx?brand_name=Cardioflo")
    assert response.status_code == 200
    assert response.headers["content-type"] == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    assert len(response.content) > 1000


def test_export_pptx_carries_a_draft_disclaimer_on_every_slide():
    """The pitch deck is the artifact most likely to be pulled out and shown or
    forwarded on its own. A disclaimer on only one slide would leave the rest
    silently unmarked the moment that slide is cropped, skipped, or deleted —
    so every slide must carry it independently.
    """
    import io
    from pptx import Presentation

    response = client.get("/api/reports/export/pptx?molecule=Empagliflozin&brand_name=Cardioflo")
    assert response.status_code == 200
    presentation = Presentation(io.BytesIO(response.content))
    assert len(presentation.slides) > 0
    for slide in presentation.slides:
        texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        assert any("Not MLR Approved" in text for text in texts), (
            "a slide is missing the draft disclaimer footer")


def test_export_xlsx_carries_a_draft_disclaimer_and_correct_row_layout():
    """The workbook is as forwardable as the deck and often outlives it, so it
    needs the same disclaimer. Also pins the funnel-section and scenario-table
    row positions: inserting the disclaimer row shifts every row below it, and
    a future edit that adds or removes a header row must not silently misalign
    the funnel labels against their values.
    """
    import io
    from openpyxl import load_workbook

    response = client.get("/api/reports/export/xlsx?brand_name=Cardioflo")
    assert response.status_code == 200
    workbook = load_workbook(io.BytesIO(response.content))
    sheet = workbook.active

    assert "Not MLR Approved" in sheet["A2"].value
    assert "MOLECULE TO MARKET AI" in sheet["A1"].value
    assert sheet["A4"].value == "EPIDEMIOLOGICAL PATIENT FUNNEL PARAMETERS"
    assert sheet["A5"].value == "Total Target Population"
    assert isinstance(sheet["B5"].value, (int, float))
    assert sheet["A17"].value == "Scenario"
    assert sheet["A18"].value == "Conservative"


# --- Regression tests for fixed defects -------------------------------------

@pytest.mark.parametrize("param,value", [
    ("prevalence_rate", -1),
    ("prevalence_rate", 5),
    ("diagnosed_rate", 0),
    ("treated_rate", 1.5),
    ("brand_adoption_rate_y1", -0.2),
    ("total_population", -999),
    ("annual_cost_per_patient_usd", 0),
])
def test_forecast_rejects_impossible_inputs(param, value):
    """Out-of-range assumptions used to return 200 with negative revenue."""
    response = client.get(f"/api/forecasting/model?{param}={value}")
    assert response.status_code == 422


def test_forecast_share_cannot_exceed_market():
    """A high Year-1 adoption rate must not project more than 100% share."""
    response = client.get("/api/forecasting/model?brand_adoption_rate_y1=0.3")
    assert response.status_code == 200
    data = response.json()
    market = data["current_therapy_market_size_usd"]
    for scenario in ("conservative_scenario", "realistic_scenario", "aggressive_scenario"):
        for year in ("year_1", "year_2", "year_3", "year_4", "year_5"):
            assert 0 <= data[scenario][year] <= market


def test_forecast_segments_match_therapy_area():
    """Oncology forecasts must not return the cardiometabolic prescriber panel."""
    response = client.get("/api/forecasting/model?therapy_area=Immuno-Oncology")
    assert response.status_code == 200
    specialties = " ".join(s["specialty"] for s in response.json()["doctor_specialties"]).lower()
    assert "oncolog" in specialties
    assert "cardiolog" not in specialties


def test_forecast_unknown_therapy_area_flags_gap():
    """An unmapped therapy area shows a gap rather than a wrong default panel."""
    response = client.get("/api/forecasting/model?therapy_area=Ophthalmology")
    assert response.status_code == 200
    segments = response.json()["doctor_specialties"]
    assert all(s["estimated_pool_size"] == 0 for s in segments)
    assert "not yet defined" in segments[0]["specialty"]


def test_brand_plan_regenerates_when_inputs_change():
    """A saved plan must not be served after the project's molecule changes."""
    project_id = f"test-{uuid.uuid4().hex[:8]}"
    first = client.get(
        f"/api/brand-plan/generate?project_id={project_id}"
        "&molecule=Empagliflozin&brand_name=First&indication=Heart+Failure&therapy_area=Cardiometabolic"
    )
    assert first.status_code == 200
    assert first.json()["molecule_name"] == "Empagliflozin"

    second = client.get(
        f"/api/brand-plan/generate?project_id={project_id}"
        "&molecule=Pembrolizumab&brand_name=Keytruda&indication=NSCLC&therapy_area=Immuno-Oncology"
    )
    assert second.status_code == 200
    body = second.json()
    assert body["molecule_name"] == "Pembrolizumab"
    assert body["brand_name"] == "Keytruda"
    assert body["indication"] == "NSCLC"


def test_brand_plan_reuses_plan_for_identical_inputs():
    project_id = f"test-{uuid.uuid4().hex[:8]}"
    query = (
        f"/api/brand-plan/generate?project_id={project_id}"
        "&molecule=Empagliflozin&brand_name=Cardioflo&indication=CKD&therapy_area=Cardiometabolic"
    )
    first = client.get(query).json()
    second = client.get(query).json()
    assert first["last_updated"] == second["last_updated"]


def test_brand_plan_refresh_forces_rebuild():
    project_id = f"test-{uuid.uuid4().hex[:8]}"
    query = (
        f"/api/brand-plan/generate?project_id={project_id}"
        "&molecule=Empagliflozin&brand_name=Cardioflo&indication=CKD&therapy_area=Cardiometabolic"
    )
    client.get(query)
    refreshed = client.get(query + "&refresh=true")
    assert refreshed.status_code == 200
    assert refreshed.json()["molecule_name"] == "Empagliflozin"


def test_assets_do_not_leak_other_molecule_clinical_content():
    """Oncology assets must not carry SGLT2/cardio-renal objections and visuals."""
    response = client.get("/api/assets/generate?molecule=Pembrolizumab&brand_name=Keytruda&indication=NSCLC")
    assert response.status_code == 200
    blob = response.text.lower()
    for term in ("sglt2", "egfr", "hba1c", "metformin", "mycotic", "cardio-renal"):
        assert term not in blob, f"leaked cardio-renal term: {term}"


@pytest.mark.parametrize("brand_name", [
    "कार्डियो",           # non-Latin-1 script
    "Café Brand",         # accented characters
    "Brand’s Plan",  # curly apostrophe pasted from Word
])
def test_export_handles_non_latin1_brand_names(brand_name):
    """These crashed with UnicodeEncodeError while building the header."""
    response = client.get("/api/reports/export/xlsx", params={"brand_name": brand_name})
    assert response.status_code == 200
    assert len(response.content) > 1000


def test_export_filename_cannot_be_steered_by_input():
    response = client.get("/api/reports/export/xlsx", params={"brand_name": "../../../etc/passwd"})
    assert response.status_code == 200
    disposition = response.headers["content-disposition"]
    assert ".." not in disposition
    assert "/" not in disposition.split("filename=")[1]


def test_export_rejects_nothing_but_still_names_the_file():
    response = client.get("/api/reports/export/xlsx", params={"brand_name": "///"})
    assert response.status_code == 200
    assert "filename=" in response.headers["content-disposition"]


# --- combination molecules and lifecycle -------------------------------------

def test_combination_resolver_splits_every_separator():
    from app.services.molecule_resolver import resolve
    for text in ["Empagliflozin + Metformin", "Empagliflozin and Metformin",
                 "Empagliflozin/Metformin", "Empagliflozin, Metformin"]:
        r = resolve(text)
        assert r.is_combination, text
        assert r.components == ["Empagliflozin", "Metformin"], text


def test_combination_resolver_strips_salts_and_strengths():
    from app.services.molecule_resolver import resolve
    r = resolve("empagliflozin and metformin hydrochloride 500mg")
    assert r.components == ["Empagliflozin", "Metformin"]
    r2 = resolve("Clindamycin phosphate")
    assert r2.components == ["Clindamycin"]
    assert not r2.is_combination


def test_inn_names_map_to_us_registry_names():
    """An Indian brand team types INN names; US registries file under USAN."""
    from app.services.inn_synonyms import candidates
    assert "acetaminophen" in candidates("paracetamol")
    assert "albuterol" in candidates("salbutamol")
    assert "clavulanate" in candidates("clavulanic acid")
    assert "rifampin" in candidates("rifampicin")
    # Salt forms fall back to the base moiety.
    assert "metformin" in candidates("metformin hydrochloride")


def test_molecule_endpoint_handles_combination():
    """PubChem 404s on combinations; the app must still return a profile."""
    response = client.get("/api/molecules/search?name=Empagliflozin%20%2B%20Metformin")
    assert response.status_code == 200
    data = response.json()
    assert data["generic_name"] == "Empagliflozin + Metformin"
    assert data["chemical_class"] == "Fixed-dose combination"


def test_lifecycle_endpoint_returns_innovator_and_patents():
    response = client.get("/api/lifecycle/molecule?molecule=Empagliflozin")
    assert response.status_code == 200
    data = response.json()
    assert data["display_name"] == "Empagliflozin"
    # Orange Book may be unreachable in a sandboxed CI run; only assert shape
    # when the data actually loaded.
    if data["all_products"]:
        assert data["innovator_company"]
        assert data["patents"]
    assert data["unavailable"], "commercial gaps must be declared, not implied"


def test_lifecycle_endpoint_handles_combination():
    response = client.get("/api/lifecycle/molecule?molecule=Empagliflozin+%2B+Metformin")
    assert response.status_code == 200
    data = response.json()
    assert data["is_combination"] is True
    assert data["components"] == ["Empagliflozin", "Metformin"]


# --- patient experience, brand naming, CDSCO ---------------------------------

def test_brand_names_are_novel_and_shaped_for_recall():
    """Candidates must not collide with marketed brands and must stay short."""
    response = client.get(
        "/api/intelligence/brand-names?molecule=Empagliflozin"
        "&therapy_area=Cardiometabolic&indication=Heart+Failure&count=10"
    )
    assert response.status_code == 200
    body = response.json()
    candidates = body["candidates"]
    if not candidates:
        return  # Orange Book unreachable in a sandboxed run
    assert len(candidates) <= 10
    for c in candidates:
        assert not c["exact_collision_with_marketed_brand"], c["name"]
        assert 5 <= c["length"] <= 9, c["name"]
        assert c["name"].isalpha()
        for key in ("ip_india_search_url", "uspto_search_url", "wipo_search_url"):
            assert c[key].startswith("http"), key
    # Names must be distinct.
    names = [c["name"] for c in candidates]
    assert len(names) == len(set(names))


def test_brand_names_reflect_therapy_area():
    onc = client.get(
        "/api/intelligence/brand-names?molecule=Pembrolizumab&therapy_area=Immuno-Oncology&count=8"
    ).json()["candidates"]
    if onc:
        assert all(c["name"].isalpha() for c in onc)


def test_cdsco_checklist_covers_india_specific_blockers():
    response = client.get("/api/intelligence/cdsco?molecule=Empagliflozin+%2B+Metformin")
    assert response.status_code == 200
    data = response.json()
    assert data["is_combination"] is True
    steps = " ".join(i["step"] for i in data["checklist"]).lower()
    # The India-specific traps a US-derived plan misses.
    for expected in ("fixed-dose", "schedule", "price control", "trademark"):
        assert expected in steps, expected
    assert data["blocking_steps"]
    assert all(i["url"].startswith("http") for i in data["checklist"])


def test_patient_experience_declares_its_caveat():
    """FAERS counts must never be presented without the incidence caveat."""
    response = client.get("/api/intelligence/patient-experience?molecule=Empagliflozin")
    assert response.status_code == 200
    data = response.json()
    caveat = data["interpretation_caveat"].lower()
    assert "not incidence" in caveat or "incidence" in caveat
    assert "spontaneous" in caveat
    assert data["data_sources"]

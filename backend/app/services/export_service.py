import io
import docx
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from typing import Dict, Any, List, Optional
from ..models.brand_plan import CompleteBrandPlan
from ..models.assets import CreativeCommercialAssets
from ..models.forecast import MarketForecast
from ..models.molecule import MoleculeProfile

def generate_brand_plan_docx(plan: CompleteBrandPlan, molecule: Optional[MoleculeProfile] = None) -> io.BytesIO:
    """Generates a professional 30+ page equivalent Word Document (.docx) for the Brand Plan."""
    doc = docx.Document()
    
    # Title Section
    title_p = doc.add_paragraph()
    title_run = title_p.add_run(f"PHARMACEUTICAL BRAND STRATEGY PLAN\n{plan.brand_name.upper()}")
    title_run.bold = True
    title_run.font.size = Pt(24)
    title_run.font.color.rgb = RGBColor(15, 76, 129) # Classic Classic Navy
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    sub_p = doc.add_paragraph()
    sub_run = sub_p.add_run(f"Target Molecule: {plan.molecule_name} | Indication: {plan.indication} | Geography: {plan.target_geography}\nDate: {plan.last_updated} | MLR Audit Status: Review Required")
    sub_run.font.size = Pt(11)
    sub_run.font.italic = True
    sub_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_paragraph("\n" + "="*50 + "\n")
    
    # Executive Vision & Mission
    doc.add_heading("1. Strategic Mission & Brand Charter", level=1)
    p_m = doc.add_paragraph()
    p_m.add_run("Brand Mission: ").bold = True
    p_m.add_run(plan.mission)
    
    p_v = doc.add_paragraph()
    p_v.add_run("Brand Vision: ").bold = True
    p_v.add_run(plan.vision)
    
    p_o = doc.add_paragraph()
    p_o.add_run("Core Commercial Objective: ").bold = True
    p_o.add_run(plan.brand_objective)
    
    p_pos = doc.add_paragraph()
    p_pos.add_run("Positioning Statement: ").bold = True
    p_pos.add_run(plan.positioning_statement)
    
    # Sections Loop
    for sec in plan.sections:
        doc.add_heading(sec.section_title, level=2)
        doc.add_paragraph(sec.content_markdown)
        if sec.key_takeaways:
            p_t = doc.add_paragraph()
            p_t.add_run("Key Strategic Takeaways:").bold = True
            for t in sec.key_takeaways:
                doc.add_paragraph(f"• {t}", style='List Bullet')
    
    # Tactical Action Plan Table
    doc.add_heading("12-Month Tactical Launch Milestones", level=1)
    table = doc.add_table(rows=1, cols=4)
    table.style = 'Table Grid'
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = 'Month'
    hdr_cells[1].text = 'Activity'
    hdr_cells[2].text = 'Responsible Team'
    hdr_cells[3].text = 'Status'
    
    for m in plan.monthly_action_plan:
        row_cells = table.add_row().cells
        row_cells[0].text = m.month_name
        row_cells[1].text = m.activity
        row_cells[2].text = m.responsible_team
        row_cells[3].text = m.status
        
    # KPI Scorecard Table
    doc.add_heading("Balanced Commercial Scorecard & KPIs", level=1)
    kpi_table = doc.add_table(rows=1, cols=5)
    kpi_table.style = 'Table Grid'
    k_hdr = kpi_table.rows[0].cells
    k_hdr[0].text = 'KPI Metric'
    k_hdr[1].text = 'Category'
    k_hdr[2].text = 'Target Q1'
    k_hdr[3].text = 'Target Q2'
    k_hdr[4].text = 'Target Year 1'
    
    for k in plan.kpi_scorecard:
        k_row = kpi_table.add_row().cells
        k_row[0].text = k.kpi_name
        k_row[1].text = k.category
        k_row[2].text = k.target_q1
        k_row[3].text = k.target_q2
        k_row[4].text = k.target_year1
        
    # Compliance Footer
    doc.add_paragraph("\n\n" + "-"*50)
    p_disc = doc.add_paragraph()
    r_disc_title = p_disc.add_run("Compliance Disclaimer: ")
    r_disc_title.bold = True
    r_disc_title.font.size = Pt(9)
    r_disc_body = p_disc.add_run("This Brand Plan is a draft for internal strategic decision-making. It is not approved promotional material. All clinical, safety, efficacy, regulatory, market, and trademark statements must undergo formal Medical-Legal-Regulatory and legal review before use.")
    r_disc_body.font.size = Pt(9)
    
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def generate_pitch_deck_pptx(
    plan: CompleteBrandPlan,
    assets: CreativeCommercialAssets,
    *,
    competitor_data: Optional[Dict[str, Any]] = None,
    regulatory: Optional[Dict[str, Any]] = None,
    molecule: Optional[Dict[str, Any]] = None,
    forecast: Optional[MarketForecast] = None,
) -> io.BytesIO:
    """Generates an executive PowerPoint pitch deck (.pptx), data-dense where real data exists.

    Every table and stat tile is built from the plan's own saved fields plus
    whatever grounding was passed in (competitor intelligence, regulatory
    status, molecule profile, market forecast) — never invented for this
    export. A section whose underlying data is missing says so explicitly
    rather than being padded with placeholder rows to look complete.
    """
    prs = Presentation()
    prs.slide_width = PptxInches(13.333) # 16:9 widescreen
    prs.slide_height = PptxInches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    def add_mlr_footer(slide) -> None:
        """Draft status on every slide.

        A pitch deck is the artifact most likely to be pulled out of the app
        and presented, forwarded, or left in an inbox on its own — unlike the
        DOCX, which carries its disclaimer as running text a reader has to
        scroll past, a single deleted or skipped slide here would leave every
        other slide silently unmarked. Repeating it as a footer means no slide
        can be extracted without carrying the warning.
        """
        footer_box = slide.shapes.add_textbox(
            PptxInches(0.8), PptxInches(7.05), PptxInches(11.7), PptxInches(0.35))
        f_tf = footer_box.text_frame
        f_tf.word_wrap = False
        p = f_tf.paragraphs[0]
        p.text = ("DRAFT — Not MLR Approved — Internal Use Only. "
                 "Clinical, safety, and efficacy statements require source verification before use.")
        p.font.size = PptxPt(9)
        p.font.italic = True
        p.font.color.rgb = PptxRGBColor(148, 163, 184)

    def add_standard_slide(title_text: str, subtitle_text: str, bullets: list):
        slide = prs.slides.add_slide(blank_layout)
        
        # Header banner
        header_box = slide.shapes.add_textbox(PptxInches(0.8), PptxInches(0.6), PptxInches(11.7), PptxInches(1.2))
        tf = header_box.text_frame
        tf.word_wrap = True
        p_title = tf.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = PptxPt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = PptxRGBColor(15, 76, 129)
        
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle_text
        p_sub.font.size = PptxPt(14)
        p_sub.font.color.rgb = PptxRGBColor(100, 110, 120)
        
        # Content box
        content_box = slide.shapes.add_textbox(PptxInches(0.8), PptxInches(2.0), PptxInches(11.7), PptxInches(4.8))
        c_tf = content_box.text_frame
        c_tf.word_wrap = True
        
        for i, b in enumerate(bullets):
            p = c_tf.add_paragraph() if i > 0 else c_tf.paragraphs[0]
            p.text = f"•  {b}"
            p.font.size = PptxPt(18)
            p.space_after = PptxPt(14)
            p.font.color.rgb = PptxRGBColor(30, 41, 59)
        
        add_mlr_footer(slide)
        return slide

    def add_table_slide(title_text: str, subtitle_text: str, headers: List[str], rows: List[List[str]], empty_note: str = "No data on file yet."):
        """A slide whose body is a table — used for anything measured (competitors,
        regulatory status, KPIs, milestones) rather than narrative bullets.
        """
        slide = prs.slides.add_slide(blank_layout)

        header_box = slide.shapes.add_textbox(PptxInches(0.6), PptxInches(0.5), PptxInches(12.1), PptxInches(1.0))
        tf = header_box.text_frame
        tf.word_wrap = True
        p_title = tf.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = PptxPt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = PptxRGBColor(15, 76, 129)
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle_text
        p_sub.font.size = PptxPt(13)
        p_sub.font.color.rgb = PptxRGBColor(100, 110, 120)

        if not rows:
            note_box = slide.shapes.add_textbox(PptxInches(0.6), PptxInches(2.0), PptxInches(12.1), PptxInches(1.0))
            note_p = note_box.text_frame.paragraphs[0]
            note_p.text = empty_note
            note_p.font.size = PptxPt(16)
            note_p.font.italic = True
            note_p.font.color.rgb = PptxRGBColor(100, 110, 120)
            add_mlr_footer(slide)
            return slide

        n_rows, n_cols = len(rows) + 1, len(headers)
        table_shape = slide.shapes.add_table(
            n_rows, n_cols, PptxInches(0.6), PptxInches(1.7), PptxInches(12.1), PptxInches(min(5.0, 0.4 * n_rows)),
        )
        table = table_shape.table
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = PptxPt(12)
            cell.text_frame.paragraphs[0].font.color.rgb = PptxRGBColor(255, 255, 255)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PptxRGBColor(15, 76, 129)
        for row_idx, row in enumerate(rows, start=1):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)
                cell.text_frame.paragraphs[0].font.size = PptxPt(11)
                cell.text_frame.paragraphs[0].font.color.rgb = PptxRGBColor(30, 41, 59)
                cell.fill.solid()
                cell.fill.fore_color.rgb = PptxRGBColor(248, 250, 252) if row_idx % 2 else PptxRGBColor(237, 242, 247)

        add_mlr_footer(slide)
        return slide

    def add_stat_dashboard_slide(title_text: str, subtitle_text: str, stats: List[tuple]):
        """Big-number callout tiles for the executive summary — each `(value, label)`."""
        slide = prs.slides.add_slide(blank_layout)
        header_box = slide.shapes.add_textbox(PptxInches(0.6), PptxInches(0.5), PptxInches(12.1), PptxInches(1.0))
        tf = header_box.text_frame
        tf.word_wrap = True
        p_title = tf.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = PptxPt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = PptxRGBColor(15, 76, 129)
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle_text
        p_sub.font.size = PptxPt(13)
        p_sub.font.color.rgb = PptxRGBColor(100, 110, 120)

        tile_width = PptxInches(3.9)
        tile_height = PptxInches(2.2)
        gap = PptxInches(0.15)
        start_x = PptxInches(0.6)
        start_y = PptxInches(2.0)
        for idx, (value, label) in enumerate(stats[:6]):
            col, row = idx % 3, idx // 3
            x = start_x + col * (tile_width + gap)
            y = start_y + row * (tile_height + gap)
            box = slide.shapes.add_textbox(x, y, tile_width, tile_height)
            tf2 = box.text_frame
            tf2.word_wrap = True
            p_val = tf2.paragraphs[0]
            p_val.text = str(value)
            p_val.font.size = PptxPt(30)
            p_val.font.bold = True
            p_val.font.color.rgb = PptxRGBColor(15, 76, 129)
            p_lab = tf2.add_paragraph()
            p_lab.text = str(label)
            p_lab.font.size = PptxPt(12)
            p_lab.font.color.rgb = PptxRGBColor(71, 85, 105)

        add_mlr_footer(slide)
        return slide

    # Slide 1: Title Slide
    s1 = prs.slides.add_slide(blank_layout)
    s1_box = s1.shapes.add_textbox(PptxInches(1.0), PptxInches(2.0), PptxInches(11.3), PptxInches(3.5))
    s1_tf = s1_box.text_frame
    s1_title = s1_tf.paragraphs[0]
    s1_title.text = f"{plan.brand_name.upper()}\nCOMMERCIAL BRAND STRATEGY"
    s1_title.font.size = PptxPt(36)
    s1_title.font.bold = True
    s1_title.font.color.rgb = PptxRGBColor(15, 76, 129)

    s1_sub = s1_tf.add_paragraph()
    s1_sub.text = f"Target Molecule: {plan.molecule_name} | {plan.therapy_area} | {plan.target_geography}\nLaunch Strategy & Commercial Operating Plan"
    s1_sub.font.size = PptxPt(18)
    s1_sub.font.color.rgb = PptxRGBColor(71, 85, 105)
    s1_sub.space_before = PptxPt(20)
    s1_dev = s1_tf.add_paragraph()
    s1_dev.text = "Molecule to Market AI — Developed by Animesh Mishra"
    s1_dev.font.size = PptxPt(11)
    s1_dev.font.italic = True
    s1_dev.font.color.rgb = PptxRGBColor(148, 163, 184)
    s1_dev.space_before = PptxPt(24)
    add_mlr_footer(s1)

    # Slide 2: Executive Summary Dashboard — real stat tiles, only for what's on file
    stats: List[tuple] = []
    if molecule and molecule.get("pharmacological_class") and molecule["pharmacological_class"] != "Not verified":
        stats.append((molecule["pharmacological_class"], "Pharmacological class"))
    market_summary = (competitor_data or {}).get("market_summary") or {}
    if market_summary.get("has_data"):
        stats.append((f"{market_summary.get('market_size')} {market_summary.get('value_unit', '')}", f"Measured market size ({market_summary.get('period', 'latest period')})"))
        stats.append((str(market_summary.get("total_brands", "—")), "Competing brands on file"))
    us_fda = (regulatory or {}).get("us_fda") or {}
    if us_fda.get("status") and us_fda["status"] != "Investigational":
        year = f", {us_fda['approval_year']}" if us_fda.get("approval_year") else ""
        stats.append((f"{us_fda['status']}{year}", "US FDA status"))
    if forecast:
        stats.append((f"{forecast.treated_patient_pool:,}", "Treated patient pool (modeled)"))
        stats.append((f"${forecast.realistic_scenario.year_5:,.0f}", "Year-5 realistic revenue (modeled)"))
    if not stats:
        stats = [("No verified data yet", "Run Modules 1, 4, 6, and 7 to populate this dashboard")]
    add_stat_dashboard_slide(
        "Executive Summary Dashboard",
        "Every figure below is read from a verified module — none are estimated for this export",
        stats,
    )

    # Slide 3: Executive Charter
    add_standard_slide(
        "Executive Brand Charter & Core Objective",
        "Defining our market ambition and clinical mission",
        [
            f"Mission: {plan.mission}",
            f"Vision: {plan.vision}",
            f"Commercial Objective: {plan.brand_objective}",
            "Target Audience: Tier A Cardiologists, Endocrinologists, Nephrologists, and Primary Care."
        ]
    )

    # Slide 4: Molecule & Regulatory Snapshot
    reg_rows: List[List[str]] = []
    if molecule:
        for key, label in (("pharmacological_class", "Pharmacological class"), ("mechanism_of_action", "Mechanism of action"), ("chemical_class", "Chemical class")):
            value = molecule.get(key)
            if value and str(value) != "Not verified":
                reg_rows.append([label, str(value)])
    for agency_key, agency_label in (("us_fda", "US FDA"), ("india_cdsco", "India CDSCO"), ("eu_ema", "EU EMA")):
        agency = (regulatory or {}).get(agency_key) or {}
        status = agency.get("status")
        if status and status != "Investigational":
            detail = status
            if agency.get("approval_year"):
                detail += f" ({agency['approval_year']})"
            app_count = len(agency.get("application_numbers") or [])
            if app_count:
                detail += f", {app_count} application(s) on file"
            reg_rows.append([agency_label, detail])
    if regulatory and regulatory.get("generic_vs_innovator_status"):
        reg_rows.append(["Market status", regulatory["generic_vs_innovator_status"]])
    if regulatory and regulatory.get("patent_expiry_timeline"):
        reg_rows.append(["Patent / exclusivity timeline", regulatory["patent_expiry_timeline"]])
    add_table_slide(
        "Molecule & Regulatory Snapshot",
        "Verified facts only — from PubChem and national regulatory labels",
        ["Field", "Verified value"], reg_rows,
        empty_note="No verified molecule or regulatory data on file yet — check Modules 1 and 4.",
    )

    # Slide 5: Competitive Landscape
    competitor_rows = [
        [c.get("brand_name", "—"), c.get("company", "—"),
         f"{c['market_share_percentage']:.1f}%" if c.get("market_share_percentage") else "—",
         "Team-attested" if c.get("data_source") == "manual" else "Measured"]
        for c in (competitor_data or {}).get("competitors") or []
    ][:10]
    add_table_slide(
        "Competitive Landscape",
        f"{market_summary.get('total_brands', 0)} brands on file" if market_summary.get("has_data") else "Measured facts only — no invented share or positioning",
        ["Brand", "Company", "Market share", "Source"], competitor_rows,
        empty_note="No competitor data on file yet — check Module 6.",
    )

    # Slide 6: SWOT
    swot = (competitor_data or {}).get("swot_analysis") or {}
    swot_cols = [swot.get(k) or [] for k in ("strengths", "weaknesses", "opportunities", "threats")]
    max_len = max((len(col) for col in swot_cols), default=0)
    swot_rows = [
        [col[i] if i < len(col) else "" for col in swot_cols]
        for i in range(max_len)
    ]
    add_table_slide(
        "SWOT Analysis",
        "Analyst-curated where available — never extended beyond the source",
        ["Strengths", "Weaknesses", "Opportunities", "Threats"], swot_rows,
        empty_note="No SWOT analysis on file yet for this molecule — check Module 6.",
    )

    # Slide 7: Strategic Positioning
    add_standard_slide(
        "Strategic Brand Positioning & RTBs",
        "Differentiating against incumbent standard of care",
        [
            f"Positioning Statement: {plan.positioning_statement}",
            f"Campaign Theme: {assets.campaign_theme}",
            f"Brand Promise & RTBs: {plan.brand_promise_and_rtb}",
            f"Competitive Gap: {plan.competitor_gap_and_differentiation}",
        ]
    )

    # Slide 8: Epidemiological Funnel & Revenue Scenarios
    forecast_rows: List[List[str]] = []
    if forecast:
        forecast_rows = [
            ["Prevalent patient pool", f"{forecast.prevalent_patient_pool:,}"],
            ["Diagnosed patient pool", f"{forecast.diagnosed_patient_pool:,}"],
            ["Treated patient pool", f"{forecast.treated_patient_pool:,}"],
            ["Year-1 revenue (conservative / realistic / aggressive)",
             f"${forecast.conservative_scenario.year_1:,.0f} / ${forecast.realistic_scenario.year_1:,.0f} / ${forecast.aggressive_scenario.year_1:,.0f}"],
            ["Year-5 revenue (conservative / realistic / aggressive)",
             f"${forecast.conservative_scenario.year_5:,.0f} / ${forecast.realistic_scenario.year_5:,.0f} / ${forecast.aggressive_scenario.year_5:,.0f}"],
            ["Realistic-scenario 5-year CAGR", f"{forecast.realistic_scenario.cagr_percentage:.1f}%"],
        ]
    add_table_slide(
        "Epidemiological Funnel & Revenue Scenarios",
        f"Modeled from a {forecast.total_population:,}-person population at {forecast.prevalence_rate*100:.1f}% prevalence" if forecast else "Modeled forecast",
        ["Metric", "Value"], forecast_rows,
        empty_note="No forecast on file yet — check Module 7.",
    )

    # Slide 9: India Trade Price Structure (only when supplied)
    if forecast and forecast.trade_price_structure:
        tps = forecast.trade_price_structure
        add_table_slide(
            "India Trade Price Structure",
            "Per patient-year, INR — manufacturer realization is PTS, not MRP",
            ["Metric", "Value"],
            [
                ["MRP (patient pays)", f"₹{tps.mrp_per_patient_year:,.0f}"],
                ["PTR (price to retailer)", f"₹{tps.ptr_per_patient_year:,.0f}"],
                ["PTS (price to stockist)", f"₹{tps.pts_per_patient_year:,.0f}"],
                ["Retailer margin", f"₹{tps.retailer_margin_amount:,.0f} ({tps.retailer_margin_percent:.1f}%)"],
                ["Stockist margin", f"₹{tps.stockist_margin_amount:,.0f} ({tps.stockist_margin_percent:.1f}%)"],
                ["Manufacturer realization of MRP", f"{tps.manufacturer_realization_percent_of_mrp:.1f}%"],
            ],
        )

    # Slide 10: Visual Aid Storyboard (Detailer Flow)
    v_bullets = [f"Slide {s.slide_number}: {s.headline_for_doctor}" for s in assets.visual_aid_slides[:4]]
    add_standard_slide(
        "Field Force Detailing Flow (Visual Aid Concept)",
        "6-Step Doctor Conversation Structure",
        v_bullets
    )

    # Slide 11: KPI Scorecard
    kpi_rows = [[k.kpi_name, k.category, k.target_q1, k.target_q2, k.target_q4, k.target_year1] for k in plan.kpi_scorecard][:10]
    add_table_slide(
        "KPI Scorecard",
        "Targets by quarter, as defined in the brand plan",
        ["KPI", "Category", "Q1", "Q2", "Q4", "Year 1"], kpi_rows,
        empty_note="No KPI scorecard on file yet — check Module 8.",
    )

    # Slide 12: 12-Month Launch Milestones
    milestone_rows = [[m.month_name, m.activity, m.responsible_team, m.status] for m in plan.monthly_action_plan][:12]
    add_table_slide(
        "Commercial Launch Milestones & Roadmap",
        "Key operational deliverables for Year 1 execution",
        ["Month", "Activity", "Owner", "Status"], milestone_rows,
        empty_note="No milestone plan on file yet — check Module 8.",
    )

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

def generate_financial_model_xlsx(forecast: MarketForecast, brand_name: str) -> io.BytesIO:
    """Generates an editable, multi-tab Excel spreadsheet (.xlsx) with forecasting formulas."""
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "5-Year Revenue Forecast"
    
    # Header styling
    header_font = Font(name="Calibri", size=14, bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="0F4C81", end_color="0F4C81", fill_type="solid")
    bold_font = Font(name="Calibri", size=11, bold=True)
    curr_format = "$#,##0"
    
    ws.merge_cells("A1:G1")
    ws["A1"] = f"MOLECULE TO MARKET AI — FINANCIAL FORECAST MODEL ({brand_name.upper()})"
    ws["A1"].font = header_font
    ws["A1"].fill = header_fill
    ws["A1"].alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 22

    # Row 2: draft status. The DOCX and PPTX both carry this; a spreadsheet is
    # just as forwardable and often outlives the deck it was built alongside,
    # so the same statement belongs on the one sheet this workbook has.
    ws.merge_cells("A2:G2")
    ws["A2"] = ("DRAFT — Not MLR Approved — Internal planning use only. "
               "Pricing, forecast, and CAGR assumptions require sourced verification before use.")
    ws["A2"].font = Font(name="Calibri", size=9, italic=True, color="94A3B8")
    ws["A2"].alignment = Alignment(horizontal="center")

    # Patient Funnel Section
    ws["A4"] = "EPIDEMIOLOGICAL PATIENT FUNNEL PARAMETERS"
    ws["A4"].font = bold_font
    
    funnel_data = [
        ("Total Target Population", forecast.total_population, "#,##0"),
        ("Disease Prevalence Rate", forecast.prevalence_rate, "0.0%"),
        ("Prevalent Patient Pool", forecast.prevalent_patient_pool, "#,##0"),
        ("Diagnosis Rate", forecast.diagnosed_rate, "0.0%"),
        ("Diagnosed Patient Pool", forecast.diagnosed_patient_pool, "#,##0"),
        ("Treatment Rate", forecast.treated_rate, "0.0%"),
        ("Treated Patient Pool", forecast.treated_patient_pool, "#,##0"),
        ("Annual Net Price Per Patient (USD)", forecast.annual_cost_per_patient_usd, "$#,##0.00"),
        ("Total Available Treated Market Size (USD)", forecast.current_therapy_market_size_usd, "$#,##0")
    ]
    
    for row_idx, (label, val, fmt) in enumerate(funnel_data, start=5):
        ws.cell(row=row_idx, column=1, value=label)
        c = ws.cell(row=row_idx, column=2, value=val)
        c.number_format = fmt
        c.alignment = Alignment(horizontal="right")
        
    # Scenario Comparison Section
    start_row = 16
    ws.cell(row=start_row, column=1, value="5-YEAR REVENUE SCENARIO PROJECTIONS (USD)").font = bold_font
    
    headers = ["Scenario", "Year 1", "Year 2", "Year 3", "Year 4", "Year 5", "5-Yr CAGR"]
    for col_idx, h in enumerate(headers, start=1):
        cell = ws.cell(row=start_row+1, column=col_idx, value=h)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(start_color="334155", end_color="334155", fill_type="solid")
        
    scenarios = [
        ("Conservative", forecast.conservative_scenario),
        ("Realistic (Base-Case)", forecast.realistic_scenario),
        ("Aggressive", forecast.aggressive_scenario)
    ]
    
    for s_idx, (s_name, s_data) in enumerate(scenarios, start=start_row+2):
        ws.cell(row=s_idx, column=1, value=s_name).font = bold_font
        for c_idx, yr_val in enumerate([s_data.year_1, s_data.year_2, s_data.year_3, s_data.year_4, s_data.year_5], start=2):
            cell = ws.cell(row=s_idx, column=c_idx, value=yr_val)
            cell.number_format = curr_format
            cell.alignment = Alignment(horizontal="right")
        cagr_cell = ws.cell(row=s_idx, column=7, value=s_data.cagr_percentage / 100.0)
        cagr_cell.number_format = "0.0%"
        cagr_cell.alignment = Alignment(horizontal="right")
        
    # Auto-adjust column widths
    for col in ws.columns:
        max_len = max(len(str(cell.value or '')) for cell in col)
        col_letter = openpyxl.utils.get_column_letter(col[0].column)
        ws.column_dimensions[col_letter].width = max(max_len + 3, 14)

    # India trade price structure — a separate sheet, not columns bolted onto
    # the first one, since PTR/PTS is an India-specific concept in INR while
    # the sheet above is USD. Only added when the forecast actually carries
    # trade pricing, so a plan with none does not show an empty tab.
    if forecast.trade_price_structure:
        _add_trade_price_sheet(wb, forecast, brand_name)

    buffer = io.BytesIO()
    wb.save(buffer)
    buffer.seek(0)
    return buffer


def _add_trade_price_sheet(wb, forecast: MarketForecast, brand_name: str) -> None:
    """India MRP -> PTR -> PTS trade margin structure, on its own sheet.

    PTS is the manufacturer's actual realization per patient-year — the
    correct basis for the company's own revenue forecast, distinct from the
    MRP the patient pays. Showing the margin structure explicitly, rather
    than just the final number, is what lets a brand manager sanity-check the
    trade terms against what the distribution channel would actually accept.
    """
    trade = forecast.trade_price_structure
    ws = wb.create_sheet("India Trade Price Structure")

    header_font = Font(name="Calibri", size=14, bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="0F4C81", end_color="0F4C81", fill_type="solid")
    bold_font = Font(name="Calibri", size=11, bold=True)
    inr_format = "₹#,##0"

    ws.merge_cells("A1:D1")
    ws["A1"] = f"INDIA TRADE PRICE STRUCTURE — {brand_name.upper()}"
    ws["A1"].font = header_font
    ws["A1"].fill = header_fill
    ws["A1"].alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 22

    ws.merge_cells("A2:D2")
    ws["A2"] = ("DRAFT — Not MLR Approved. Trade terms require sourced verification "
               "against the actual distribution agreement before use.")
    ws["A2"].font = Font(name="Calibri", size=9, italic=True, color="94A3B8")
    ws["A2"].alignment = Alignment(horizontal="center")

    ws["A4"] = "PRICE POINT (PER PATIENT-YEAR)"
    ws["A4"].font = bold_font
    ws["B4"] = "VALUE (INR)"
    ws["B4"].font = bold_font

    rows = [
        ("MRP — Maximum Retail Price (what the patient pays)", trade.mrp_per_patient_year),
        ("PTR — Price to Retailer", trade.ptr_per_patient_year),
        ("PTS — Price to Stockist (manufacturer's own realization)", trade.pts_per_patient_year),
    ]
    for idx, (label, value) in enumerate(rows, start=5):
        ws.cell(row=idx, column=1, value=label)
        cell = ws.cell(row=idx, column=2, value=value)
        cell.number_format = inr_format
        cell.alignment = Alignment(horizontal="right")

    ws["A9"] = "MARGIN"
    ws["A9"].font = bold_font
    ws["B9"] = "AMOUNT (INR)"
    ws["B9"].font = bold_font
    ws["C9"] = "% OF UPSTREAM PRICE"
    ws["C9"].font = bold_font

    margin_rows = [
        ("Retailer margin (MRP - PTR)", trade.retailer_margin_amount, trade.retailer_margin_percent),
        ("Stockist margin (PTR - PTS)", trade.stockist_margin_amount, trade.stockist_margin_percent),
    ]
    for idx, (label, amount, percent) in enumerate(margin_rows, start=10):
        ws.cell(row=idx, column=1, value=label)
        amount_cell = ws.cell(row=idx, column=2, value=amount)
        amount_cell.number_format = inr_format
        amount_cell.alignment = Alignment(horizontal="right")
        percent_cell = ws.cell(row=idx, column=3, value=percent / 100.0)
        percent_cell.number_format = "0.0%"
        percent_cell.alignment = Alignment(horizontal="right")

    ws["A13"] = "Manufacturer realization as % of MRP"
    ws["A13"].font = bold_font
    pct_cell = ws.cell(row=13, column=2, value=trade.manufacturer_realization_percent_of_mrp / 100.0)
    pct_cell.number_format = "0.0%"
    pct_cell.font = bold_font
    pct_cell.alignment = Alignment(horizontal="right")

    if forecast.therapy_market_size_inr_at_trade_price:
        ws["A15"] = "Addressable market at PTS (treated pool x PTS)"
        ws["A15"].font = bold_font
        market_cell = ws.cell(row=15, column=2, value=forecast.therapy_market_size_inr_at_trade_price)
        market_cell.number_format = "₹#,##0"
        market_cell.font = bold_font
        market_cell.alignment = Alignment(horizontal="right")
        ws["A16"] = ("This is the manufacturer's own addressable revenue, not the "
                     "patient-facing market size shown on the USD sheet — the two "
                     "are not meant to reconcile, since PTS excludes the retailer "
                     "and stockist margins MRP includes.")
        ws["A16"].font = Font(name="Calibri", size=9, italic=True, color="64748B")
        ws.merge_cells("A16:D16")

    for col, width in (("A", 55), ("B", 18), ("C", 20), ("D", 14)):
        ws.column_dimensions[col].width = width
